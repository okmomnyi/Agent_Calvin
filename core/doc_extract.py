"""Document text extraction + chunking for the Study Vault (Phase 9).

Extracts text from PDFs (per page), PPTX (per slide), DOCX, plain text/markdown, and
note images (OCR via tesseract, if installed) — each returned with a location label so
answers can cite file + page/slide. Chunking targets ~800 tokens with overlap. All
heavy parsers import lazily; an unsupported/failed file is skipped gracefully, never fatal.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from core.logging_setup import get_logger

log = get_logger("core.doc_extract")

SUPPORTED = {".pdf", ".pptx", ".docx", ".txt", ".md", ".png", ".jpg", ".jpeg"}


@dataclass
class Passage:
    """A located block of extracted text (before chunking)."""

    loc: str      # e.g. "p.3", "slide 2", "" for whole-doc formats
    text: str


def extract(path: str | Path) -> list[Passage]:
    """Extract located passages from a document. Returns [] for unsupported/failed files."""
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            return _pdf(p)
        if ext == ".pptx":
            return _pptx(p)
        if ext == ".docx":
            return _docx(p)
        if ext in (".txt", ".md"):
            return [Passage("", p.read_text(encoding="utf-8", errors="ignore"))]
        if ext in (".png", ".jpg", ".jpeg"):
            return _image(p)
    except Exception:  # noqa: BLE001 - one bad file must not break ingestion
        log.exception("extract failed for %s", p)
    return []


def _pdf(p: Path) -> list[Passage]:
    from pypdf import PdfReader

    reader = PdfReader(str(p))
    out = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            out.append(Passage(f"p.{i}", text))
    return out


def _pptx(p: Path) -> list[Passage]:
    from pptx import Presentation

    prs = Presentation(str(p))
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        if parts:
            out.append(Passage(f"slide {i}", "\n".join(parts)))
    return out


def _docx(p: Path) -> list[Passage]:
    import docx

    doc = docx.Document(str(p))
    text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    return [Passage("", text)] if text else []


def _image(p: Path) -> list[Passage]:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        log.warning("OCR skipped for %s (install pytesseract + tesseract binary)", p.name)
        return []
    try:
        text = pytesseract.image_to_string(Image.open(str(p))).strip()
    except Exception:  # noqa: BLE001 - tesseract binary missing / bad image
        log.warning("OCR failed for %s (is the tesseract binary installed?)", p.name)
        return []
    return [Passage("", text)] if text else []


def chunk_passages(passages: list[Passage], size_words: int = 500, overlap: int = 60) -> list[Passage]:
    """Split each passage into ~size_words chunks (with overlap), preserving its loc label."""
    chunks: list[Passage] = []
    for passage in passages:
        words = passage.text.split()
        if not words:
            continue
        if len(words) <= size_words:
            chunks.append(Passage(passage.loc, passage.text.strip()))
            continue
        start = 0
        while start < len(words):
            window = words[start:start + size_words]
            chunks.append(Passage(passage.loc, " ".join(window)))
            if start + size_words >= len(words):
                break
            start += size_words - overlap
    return chunks


# ------------------------------------------------------------------ PDF text repair
# A single capital stranded from its word by PDF kerning: "T echnology" -> "Technology".
# "I" and "A" are excluded because they are real one-letter words.
_SPLIT_WORD = re.compile(r"\b(?![IA]\b)([A-Z]) ([a-z]{2,})")
# A location glued straight onto the preceding word: "TechnologyMeru, KE".
_GLUED_PLACE = re.compile(r"([a-z])([A-Z][a-z]+,\s*[A-Z]{2}\b)")
# A date range glued onto a role: "Cloud Projects2024 - Present".
_GLUED_DATE = re.compile(r"([A-Za-z])((?:[A-Z][a-z]{2,8}\.?\s*)?\d{4}\s*[-–—])")


def clean_pdf_artifacts(text: str) -> str:
    """Repair the text damage PDF extraction does to CV-style layouts.

    Two-column CVs (institution left, location right) come out of pypdf with the columns
    concatenated and words broken by kerning:

        "Meru University of Science and T echnologyMeru, KE"
        "Independent Cybersecurity & Cloud Projects2024 - Present"

    Fed to a parser that is (correctly) told not to infer anything, this produced a CV fact
    claiming Calvin's degree was from Computer Pride in June-July 2024 -- his CompTIA course --
    when his BSc is Meru University, Sept 2024 - Apr 2028. Two entries merged into one wrong
    one, on a document employers read.

    Only separators and kerning are touched: no word is added, removed or reordered.
    """
    if not text:
        return text
    out = _SPLIT_WORD.sub(r"\1\2", text)
    out = _GLUED_PLACE.sub(r"\1  |  \2", out)
    out = _GLUED_DATE.sub(r"\1  |  \2", out)
    return out
